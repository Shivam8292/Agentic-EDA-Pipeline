"""
Export Service — PDF, PowerPoint, and Excel exports.

Uses:
- ReportLab → PDF
- python-pptx → PowerPoint
- OpenPyXL → Excel (cleaned data + stats)
"""

import io
import json
import logging
import os
from datetime import datetime
from typing import Any

import numpy as np
import pandas as pd
import plotly.io as pio
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    Image as RLImage,
    PageBreak,
    Table,
    TableStyle,
)

logger = logging.getLogger(__name__)

# Dark theme color helpers
DARK_BG = (10 / 255, 10 / 255, 15 / 255)          # #0A0A0F
CARD_BG = (17 / 255, 17 / 255, 24 / 255)           # #111118
ACCENT = (99 / 255, 102 / 255, 241 / 255)          # #6366F1
TEXT_PRIMARY = (240 / 255, 240 / 255, 245 / 255)   # #F0F0F5


def _chart_to_png_bytes(chart_json: str, width: int = 900, height: int = 500) -> bytes | None:
    """Convert a Plotly chart JSON to PNG bytes using kaleido."""
    try:
        import plotly.graph_objects as go
        fig = go.Figure(json.loads(chart_json))
        img_bytes = pio.to_image(fig, format="png", width=width, height=height, scale=1.5)
        return img_bytes
    except Exception as e:
        logger.error(f"Chart to PNG conversion failed: {e}")
        return None


# ─────────────────────────────────────────────
# PDF Export
# ─────────────────────────────────────────────

def generate_pdf(
    filename: str,
    analysis_results: list[dict[str, Any]],
    cleaning_report: dict[str, Any],
) -> bytes:
    """Generate a professional PDF report."""
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=letter,
        leftMargin=0.75 * inch,
        rightMargin=0.75 * inch,
        topMargin=0.75 * inch,
        bottomMargin=0.75 * inch,
    )

    styles = getSampleStyleSheet()
    accent_color = colors.HexColor("#6366F1")
    text_color = colors.HexColor("#0A0A0F")

    title_style = ParagraphStyle(
        "Title",
        parent=styles["Title"],
        textColor=colors.white,
        backColor=colors.HexColor("#0A0A0F"),
        fontSize=28,
        leading=34,
        spaceAfter=12,
    )
    heading_style = ParagraphStyle(
        "Heading",
        parent=styles["Heading1"],
        textColor=accent_color,
        fontSize=14,
        spaceBefore=16,
        spaceAfter=8,
    )
    body_style = ParagraphStyle(
        "Body",
        parent=styles["Normal"],
        textColor=text_color,
        fontSize=10,
        leading=15,
    )
    meta_style = ParagraphStyle(
        "Meta",
        parent=styles["Normal"],
        textColor=colors.HexColor("#4A4A5E"),
        fontSize=9,
    )

    story = []

    # Cover page
    story.append(Spacer(1, 1 * inch))
    story.append(Paragraph("DataSense Analysis Report", heading_style))
    story.append(Paragraph(f"Dataset: {filename}", body_style))
    story.append(Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y at %H:%M')}", meta_style))
    story.append(Spacer(1, 0.3 * inch))

    # Cleaning summary
    cr = cleaning_report
    cleaning_text = (
        f"Data Cleaning Summary: "
        f"{cr.get('nulls_filled', 0)} nulls filled · "
        f"{cr.get('duplicates_removed', 0)} duplicates removed · "
        f"{len(cr.get('type_corrections', []))} type corrections · "
        f"{cr.get('outliers_flagged', 0)} outliers flagged"
    )
    story.append(Paragraph(cleaning_text, meta_style))
    story.append(PageBreak())

    # Question sections
    for i, result in enumerate(analysis_results):
        if result.get("status") == "failed":
            story.append(Paragraph(f"Q{i+1}: {result['question']}", heading_style))
            story.append(Paragraph("⚠ This question could not be answered from the available data.", body_style))
            story.append(Spacer(1, 0.3 * inch))
            continue

        story.append(Paragraph(f"Q{i+1}: {result['question']}", heading_style))

        # Chart image
        chart_json = result.get("chart_json")
        if chart_json:
            img_bytes = _chart_to_png_bytes(chart_json, width=850, height=450)
            if img_bytes:
                img_buffer = io.BytesIO(img_bytes)
                rl_img = RLImage(img_buffer, width=6.5 * inch, height=3.5 * inch)
                story.append(rl_img)

        # Insight
        insight = result.get("insight", "")
        if insight:
            story.append(Spacer(1, 0.2 * inch))
            story.append(Paragraph("📊 Analyst Insight", heading_style))
            story.append(Paragraph(insight, body_style))

        story.append(Spacer(1, 0.3 * inch))
        if i < len(analysis_results) - 1:
            story.append(PageBreak())

    doc.build(story)
    return buffer.getvalue()


# ─────────────────────────────────────────────
# PowerPoint Export
# ─────────────────────────────────────────────

def generate_ppt(
    filename: str,
    analysis_results: list[dict[str, Any]],
) -> bytes:
    """Generate a dark-themed PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    dark_bg_color = RGBColor(0x0A, 0x0A, 0x0F)
    accent_color = RGBColor(0x63, 0x66, 0xF1)
    text_color = RGBColor(0xF0, 0xF0, 0xF5)
    muted_color = RGBColor(0x8B, 0x8B, 0x9E)

    def _set_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = dark_bg_color

    def _add_text_box(slide, text, left, top, width, height, font_size=18, bold=False, color=None):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        run = p.runs[0] if p.runs else p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color or text_color
        return txBox

    # Title slide
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide)

    _add_text_box(
        slide, "DataSense Analysis Report",
        Inches(1), Inches(2), Inches(11), Inches(1.5),
        font_size=40, bold=True, color=text_color
    )
    _add_text_box(
        slide, f"Dataset: {filename}",
        Inches(1), Inches(3.8), Inches(11), Inches(0.8),
        font_size=20, color=muted_color
    )
    _add_text_box(
        slide, f"Generated: {datetime.now().strftime('%B %d, %Y')}",
        Inches(1), Inches(4.7), Inches(11), Inches(0.5),
        font_size=14, color=muted_color
    )

    # Question slides
    for i, result in enumerate(analysis_results):
        slide = prs.slides.add_slide(slide_layout)
        _set_slide_bg(slide)

        # Question title
        _add_text_box(
            slide, f"Q{i+1}: {result['question']}",
            Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.8),
            font_size=18, bold=True, color=text_color
        )

        if result.get("status") == "failed":
            _add_text_box(
                slide, "⚠ This question could not be answered from the available data.",
                Inches(0.5), Inches(2), Inches(12), Inches(1),
                font_size=16, color=RGBColor(0xEF, 0x44, 0x44)
            )
            continue

        # Chart image
        chart_json = result.get("chart_json")
        if chart_json:
            img_bytes = _chart_to_png_bytes(chart_json, width=900, height=450)
            if img_bytes:
                img_stream = io.BytesIO(img_bytes)
                slide.shapes.add_picture(
                    img_stream,
                    Inches(0.5), Inches(1.2),
                    Inches(8.5), Inches(4.5)
                )

        # Insight panel (right side)
        insight = result.get("insight", "")
        if insight:
            _add_text_box(
                slide, "Analyst Insight",
                Inches(9.3), Inches(1.2), Inches(3.5), Inches(0.5),
                font_size=12, bold=True, color=accent_color
            )
            # Truncate insight to avoid overflow
            short_insight = insight[:400] + "..." if len(insight) > 400 else insight
            _add_text_box(
                slide, short_insight,
                Inches(9.3), Inches(1.8), Inches(3.5), Inches(4),
                font_size=10, color=muted_color
            )

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


# ─────────────────────────────────────────────
# Excel Export
# ─────────────────────────────────────────────

def generate_excel(
    cleaned_df: pd.DataFrame,
    cleaning_report: dict[str, Any],
) -> bytes:
    """Generate an Excel file with cleaned data, cleaning report, and column stats."""
    buffer = io.BytesIO()
    wb = Workbook()

    # ── Sheet 1: Cleaned Data ──────────────────
    ws1 = wb.active
    ws1.title = "Cleaned Data"

    header_fill = PatternFill(start_color="0A0A0F", end_color="0A0A0F", fill_type="solid")
    header_font = Font(color="F0F0F5", bold=True)
    accent_fill = PatternFill(start_color="6366F1", end_color="6366F1", fill_type="solid")

    for col_idx, col_name in enumerate(cleaned_df.columns, start=1):
        cell = ws1.cell(row=1, column=col_idx, value=col_name)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center")

    for row_idx, row in enumerate(cleaned_df.values, start=2):
        for col_idx, value in enumerate(row, start=1):
            ws1.cell(row=row_idx, column=col_idx, value=value if not isinstance(value, float) or not np.isnan(value) else None)

    # Auto-size columns
    for col in ws1.columns:
        max_length = max((len(str(cell.value or "")) for cell in col), default=8)
        ws1.column_dimensions[col[0].column_letter].width = min(max_length + 4, 40)

    # ── Sheet 2: Cleaning Report ───────────────
    ws2 = wb.create_sheet("Cleaning Report")
    ws2["A1"] = "Metric"
    ws2["B1"] = "Value"
    ws2["A1"].font = Font(bold=True, color="6366F1")
    ws2["B1"].font = Font(bold=True, color="6366F1")

    report_rows = [
        ("Nulls Filled", cleaning_report.get("nulls_filled", 0)),
        ("Duplicates Removed", cleaning_report.get("duplicates_removed", 0)),
        ("Outliers Flagged (IQR)", cleaning_report.get("outliers_flagged", 0)),
        ("Type Corrections", len(cleaning_report.get("type_corrections", []))),
    ]
    for r, (metric, value) in enumerate(report_rows, start=2):
        ws2.cell(row=r, column=1, value=metric)
        ws2.cell(row=r, column=2, value=value)

    if cleaning_report.get("type_corrections"):
        ws2.cell(row=len(report_rows) + 3, column=1, value="Type Correction Details").font = Font(bold=True)
        for i, correction in enumerate(cleaning_report["type_corrections"], start=len(report_rows) + 4):
            ws2.cell(row=i, column=1, value=correction)

    ws2.column_dimensions["A"].width = 30
    ws2.column_dimensions["B"].width = 20

    # ── Sheet 3: Column Stats ──────────────────
    ws3 = wb.create_sheet("Column Stats")
    stat_headers = ["Column", "dtype", "count", "nulls", "null_%", "mean", "median", "std", "min", "max", "unique"]
    for col_idx, h in enumerate(stat_headers, start=1):
        cell = ws3.cell(row=1, column=col_idx, value=h)
        cell.font = Font(bold=True, color="6366F1")

    for row_idx, col_name in enumerate(cleaned_df.columns, start=2):
        series = cleaned_df[col_name]
        null_count = int(series.isna().sum())
        null_pct = round(null_count / len(cleaned_df) * 100, 2)
        stats = [
            col_name,
            str(series.dtype),
            len(series),
            null_count,
            null_pct,
            round(series.mean(), 4) if pd.api.types.is_numeric_dtype(series) else "N/A",
            round(series.median(), 4) if pd.api.types.is_numeric_dtype(series) else "N/A",
            round(series.std(), 4) if pd.api.types.is_numeric_dtype(series) else "N/A",
            series.min() if pd.api.types.is_numeric_dtype(series) else "N/A",
            series.max() if pd.api.types.is_numeric_dtype(series) else "N/A",
            series.nunique(),
        ]
        for col_idx, val in enumerate(stats, start=1):
            ws3.cell(row=row_idx, column=col_idx, value=val)

    for col in ws3.columns:
        max_length = max((len(str(cell.value or "")) for cell in col), default=8)
        ws3.column_dimensions[col[0].column_letter].width = min(max_length + 4, 30)

    wb.save(buffer)
    return buffer.getvalue()
